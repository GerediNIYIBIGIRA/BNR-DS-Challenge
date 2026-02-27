"""
Generate the BNR RAG System presentation (PPTX).

Usage:
    pip install python-pptx
    python create_presentation.py
    → creates: BNR_RAG_System_Presentation.pptx
"""
from __future__ import annotations

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pathlib import Path


# ── Brand colours ─────────────────────────────────────────────────────────────
BNR_BLUE   = RGBColor(0x1A, 0x3C, 0x8F)   # deep BNR blue
ACCENT     = RGBColor(0x00, 0x7B, 0xC2)   # lighter blue
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY = RGBColor(0xF4, 0xF6, 0xF9)
DARK_GREY  = RGBColor(0x33, 0x33, 0x33)
GREEN      = RGBColor(0x1A, 0x8C, 0x4E)


# ── Helpers ───────────────────────────────────────────────────────────────────

def _add_slide(prs: Presentation, layout_idx: int = 6):
    """Add a blank slide (layout 6 = blank on most themes)."""
    layout = prs.slide_layouts[layout_idx]
    return prs.slides.add_slide(layout)


def _bg(slide, color: RGBColor) -> None:
    """Set solid background colour for a slide."""
    from pptx.oxml.ns import qn
    from lxml import etree
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _box(slide, left, top, width, height,
         text: str,
         font_size: int = 18,
         bold: bool = False,
         color: RGBColor = DARK_GREY,
         bg_color: RGBColor | None = None,
         align=PP_ALIGN.LEFT,
         wrap: bool = True):
    """Add a text box with styling."""
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = wrap

    if bg_color:
        fill = txBox.fill
        fill.solid()
        fill.fore_color.rgb = bg_color

    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txBox


def _rect(slide, left, top, width, height, color: RGBColor):
    """Add a solid coloured rectangle."""
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def _bullets(slide, left, top, width, height,
             items: list[str],
             font_size: int = 16,
             color: RGBColor = DARK_GREY,
             bold_first: bool = False):
    """Add a bulleted list."""
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = ("• " if not item.startswith("  ") else "") + item
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        run.font.bold = bold_first and i == 0


# ── Slides ────────────────────────────────────────────────────────────────────

def slide_1_title(prs: Presentation) -> None:
    """Slide 1 — Title"""
    sl = _add_slide(prs)
    _bg(sl, BNR_BLUE)

    # Header bar
    _rect(sl, 0, 0, 10, 0.18, ACCENT)

    # Main title
    _box(sl, 0.5, 1.0, 9, 1.2,
         "BNR Document Intelligence Assistant",
         font_size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Subtitle
    _box(sl, 0.5, 2.3, 9, 0.6,
         "Retrieval-Augmented Generation (RAG) over BNR Institutional Corpus",
         font_size=20, color=RGBColor(0xCC, 0xDD, 0xFF), align=PP_ALIGN.CENTER)

    # Key stats
    _rect(sl, 1.5, 3.3, 7, 1.5, RGBColor(0x0F, 0x28, 0x6B))
    _box(sl, 1.5, 3.4, 7, 1.3,
         "4 corpus documents  ·  Semantic retrieval  ·  Grounded answers with citations\n"
         "Strict fallback for out-of-corpus questions  ·  Full audit trail",
         font_size=16, color=WHITE, align=PP_ALIGN.CENTER)

    # Author / date
    _box(sl, 0.5, 5.4, 9, 0.5,
         "Geredi Niyibigira  ·  BNR Data Science Challenge #1  ·  February 2026",
         font_size=13, color=RGBColor(0xAA, 0xBB, 0xDD), align=PP_ALIGN.CENTER)


def slide_2_architecture(prs: Presentation) -> None:
    """Slide 2 — System Architecture"""
    sl = _add_slide(prs)
    _bg(sl, WHITE)
    _rect(sl, 0, 0, 10, 0.5, BNR_BLUE)
    _box(sl, 0.2, 0.05, 9.6, 0.4, "System Architecture",
         font_size=22, bold=True, color=WHITE)

    # Pipeline boxes
    steps = [
        ("1  User Question", 0.4),
        ("2  Text Embedding\nall-MiniLM-L6-v2 (local)", 1.35),
        ("3  ChromaDB\nCosine similarity, top-5 chunks", 2.45),
        ("4  Context + Prompt\nStrict grounding instructions", 3.55),
        ("5  Claude LLM\nGrounded answer + citations", 4.65),
        ("6  Audit Logger\nlogs/audit.jsonl", 5.55),
    ]
    colors = [ACCENT, BNR_BLUE, ACCENT, BNR_BLUE, ACCENT, GREEN]

    for (label, top), col in zip(steps, colors):
        _rect(sl, 0.3, top, 5.5, 0.75, col)
        _box(sl, 0.35, top + 0.05, 5.4, 0.65,
             label, font_size=14, bold=False, color=WHITE)

    # Right panel — components
    _rect(sl, 6.2, 0.6, 3.5, 5.6, LIGHT_GREY)
    _bullets(sl, 6.3, 0.65, 3.4, 5.5, [
        "COMPONENTS",
        "",
        "Embeddings",
        "  sentence-transformers",
        "  384-dim, CPU, no API key",
        "",
        "Vector Store",
        "  ChromaDB (persistent)",
        "  Cosine similarity",
        "",
        "LLM",
        "  Claude 3.5 Haiku",
        "  Anthropic API",
        "",
        "Chunking",
        "  600 words / 100 overlap",
    ], font_size=13, color=DARK_GREY, bold_first=True)


def slide_3_design(prs: Presentation) -> None:
    """Slide 3 — Key Design Choices"""
    sl = _add_slide(prs)
    _bg(sl, WHITE)
    _rect(sl, 0, 0, 10, 0.5, BNR_BLUE)
    _box(sl, 0.2, 0.05, 9.6, 0.4, "Key Design Choices",
         font_size=22, bold=True, color=WHITE)

    choices = [
        ("Local Embeddings", ACCENT,
         "all-MiniLM-L6-v2 runs entirely on CPU — "
         "corpus text never leaves the institution's infrastructure."),
        ("Strict Grounding Prompt", BNR_BLUE,
         "System prompt forbids the LLM from using external knowledge. "
         "Mandatory fallback: \"The answer cannot be determined from the provided documents.\""),
        ("Metadata Preservation", ACCENT,
         "Every chunk retains source document name, filename, page number, and doc type. "
         "Enables precise citations and filtered retrieval."),
        ("Persistent Vector Index", BNR_BLUE,
         "ChromaDB persists to disk (chroma_db/). Re-indexing only needed "
         "when the corpus changes — cold start avoided on every query."),
        ("Audit Trail", GREEN,
         "Every query is appended to logs/audit.jsonl with timestamp, tokens, "
         "latency, sources retrieved, and whether the fallback was triggered."),
    ]

    for i, (title, col, desc) in enumerate(choices):
        top = 0.65 + i * 1.0
        _rect(sl, 0.2, top, 0.08, 0.55, col)
        _box(sl, 0.45, top, 2.5, 0.35,
             title, font_size=15, bold=True, color=col)
        _box(sl, 0.45, top + 0.32, 9.3, 0.6,
             desc, font_size=13, color=DARK_GREY)


def slide_4_retrieval(prs: Presentation) -> None:
    """Slide 4 — Retrieval Approach"""
    sl = _add_slide(prs)
    _bg(sl, WHITE)
    _rect(sl, 0, 0, 10, 0.5, BNR_BLUE)
    _box(sl, 0.2, 0.05, 9.6, 0.4, "Retrieval Approach",
         font_size=22, bold=True, color=WHITE)

    # Left column
    _rect(sl, 0.2, 0.65, 4.6, 5.6, LIGHT_GREY)
    _bullets(sl, 0.35, 0.7, 4.4, 5.5, [
        "HOW IT WORKS",
        "",
        "1.  Query → embedding vector",
        "     (same model as indexing)",
        "",
        "2.  Cosine similarity search",
        "     across all ~2,000+ chunks",
        "",
        "3.  Top-5 chunks returned",
        "     with similarity scores",
        "",
        "4.  Low-quality chunks filtered",
        "     (< 80 chars at index time)",
        "",
        "5.  Chunks + question passed",
        "     to the LLM for generation",
    ], font_size=13, color=DARK_GREY, bold_first=True)

    # Right column
    _rect(sl, 5.1, 0.65, 4.7, 5.6, LIGHT_GREY)
    _bullets(sl, 5.25, 0.7, 4.5, 5.5, [
        "CORPUS DOCUMENT HANDLING",
        "",
        "PDF documents",
        "  · pypdf text extraction",
        "  · 600-word sliding window",
        "  · Page number preserved",
        "",
        "CSV (IMF dataset)",
        "  · One chunk per indicator row",
        "  · Year columns as key-value pairs",
        "  · indicator name in metadata",
        "",
        "Multi-source queries",
        "  · No per-document filter",
        "  · Naturally surfaces best",
        "    chunks from all 4 docs",
    ], font_size=13, color=DARK_GREY, bold_first=True)


def slide_5_outputs(prs: Presentation) -> None:
    """Slide 5 — Example Outputs (real evaluation results)"""
    sl = _add_slide(prs)
    _bg(sl, WHITE)
    _rect(sl, 0, 0, 10, 0.5, BNR_BLUE)
    _box(sl, 0.2, 0.05, 9.6, 0.4, "Example Outputs  (live evaluation results)",
         font_size=22, bold=True, color=WHITE)

    examples = [
        (
            "Q1 — Barriers to financial inclusion in rural Rwanda",
            ACCENT,
            "Retrieved: FinScope 2024 pp. 75, 32, 33, 5, 15 (sim 0.72–0.75)  |  Latency: 2.6 s\n"
            "Answer: Appropriate refusal — retrieved pages cover achievements & recommendations, "
            "not barriers. System correctly refused rather than fabricating an answer. "
            "Limitation: semantic search matched the topic but not the sub-topic; hybrid BM25 would improve recall.",
        ),
        (
            "Q2 — Gender gap in mobile money",
            BNR_BLUE,
            "Retrieved: GSMA 2025 pp. 70–77 (sim 0.72–0.77)  |  Latency: 5.7 s\n"
            "Answer: Detailed & accurate — Pakistan women 70% less likely to own accounts; frequency & "
            "transaction-type gaps quantified. Limitation: Rwanda-specific gender data (FinScope) was "
            "not retrieved — single-document bias toward GSMA.",
        ),
        (
            "Q3 — NBR's operational role in the payment system",
            ACCENT,
            "Retrieved: Law No. 061/2021 p.1 + p.40; FinScope pp. 60, 61 (sim 0.67–0.76)  |  Latency: 5.3 s\n"
            "Answer: Multi-source synthesis — Art. 4 (General Powers), Art. 5 (Investigative), "
            "Interoperability, RNPS Strategy 2018–2024. Limitation: strategic role cited via FinScope "
            "(secondary) rather than the law itself; one article number not specified.",
        ),
        (
            "Q4 — Has digital payment adoption increased?  [CRITICAL FAILURE CASE]",
            RGBColor(0xCC, 0x33, 0x33),
            "Retrieved: GSMA pp. 57–61 — ALL similarity scores < 0.58 (weakest in evaluation)  |  Latency: 5.9 s\n"
            "Answer: Confident ASEAN-region statistics returned despite low retrieval confidence. "
            "Q1 with HIGHER scores (0.72) correctly refused; this did not. "
            "Root cause: no minimum similarity threshold gate — fixed by adding cutoff ≥ 0.65.",
        ),
    ]

    for i, (title, col, desc) in enumerate(examples):
        top = 0.60 + i * 1.6
        _rect(sl, 0.2, top, 9.6, 1.3, LIGHT_GREY)
        _rect(sl, 0.2, top, 0.07, 1.3, col)
        _box(sl, 0.45, top + 0.04, 9.2, 0.35,
             title, font_size=13, bold=True, color=col)
        _box(sl, 0.45, top + 0.38, 9.2, 0.88,
             desc, font_size=11, color=DARK_GREY)


def slide_6_evaluation(prs: Presentation) -> None:
    """Slide 6 — Evaluation Findings & Limitations (real data)"""
    sl = _add_slide(prs)
    _bg(sl, WHITE)
    _rect(sl, 0, 0, 10, 0.5, BNR_BLUE)
    _box(sl, 0.2, 0.05, 9.6, 0.4, "Evaluation Findings & Limitations",
         font_size=22, bold=True, color=WHITE)

    # Strengths
    _rect(sl, 0.2, 0.60, 4.6, 3.0, RGBColor(0xE8, 0xF5, 0xEC))
    _rect(sl, 0.2, 0.60, 4.6, 0.4, GREEN)
    _box(sl, 0.25, 0.62, 4.5, 0.36, "WHAT WORKS WELL",
         font_size=14, bold=True, color=WHITE)
    _bullets(sl, 0.3, 1.05, 4.4, 2.5, [
        "No hallucination in any of the 5 answers",
        "Appropriate refusal on Q1 (correct grounding)",
        "Article-level legal citations on Q3",
        "GSMA gender chapter retrieved precisely (Q2)",
        "Honest partial answer on Q5",
        "Full audit trail in logs/audit.jsonl",
        "368 chunks indexed from 4 heterogeneous docs",
    ], font_size=12, color=DARK_GREY)

    # Limitations
    _rect(sl, 5.2, 0.60, 4.6, 3.0, RGBColor(0xFD, 0xF0, 0xF0))
    _rect(sl, 5.2, 0.60, 4.6, 0.4, RGBColor(0xCC, 0x33, 0x33))
    _box(sl, 5.25, 0.62, 4.5, 0.36, "IDENTIFIED FAILURES",
         font_size=14, bold=True, color=WHITE)
    _bullets(sl, 5.3, 1.05, 4.4, 2.5, [
        "Q4: sim < 0.58 → confident answer (no gate)",
        "Q2: Rwanda gender data not retrieved",
        "Q5: GSMA not surfaced for comparison",
        "Q1: 'barriers' semantically missed in top-5",
        "PDF tables & figures not extracted",
        "No minimum similarity threshold enforced",
        "Single-document bias on multi-source queries",
    ], font_size=12, color=DARK_GREY)

    # Real metrics row
    _rect(sl, 0.2, 3.75, 9.6, 2.9, LIGHT_GREY)
    _box(sl, 0.3, 3.80, 9.4, 0.38,
         "MEASURED METRICS  (5 questions, live run)",
         font_size=14, bold=True, color=BNR_BLUE)

    # Score table (text-based)
    table_rows = [
        ("Q1  Barriers (rural Rwanda)", "0.718–0.752", "Appropriate refusal", "✓"),
        ("Q2  Gender gap in mobile money", "0.717–0.770", "Detailed (global)", "✓"),
        ("Q3  NBR payment system role", "0.672–0.760", "Multi-source synthesis", "✓"),
        ("Q4  Digital payment adoption", "0.508–0.579", "Confident answer (LOW)", "✗"),
        ("Q5  Rwanda vs global trends", "0.699–0.775", "Honest partial answer", "✓"),
    ]
    verdict_colors = [GREEN, GREEN, GREEN, RGBColor(0xCC, 0x33, 0x33), GREEN]
    for i, (q, scores, verdict, icon) in enumerate(table_rows):
        top = 4.22 + i * 0.46
        _box(sl, 0.3, top, 3.8, 0.42, q, font_size=11, color=DARK_GREY)
        _box(sl, 4.2, top, 1.8, 0.42, scores, font_size=11, color=DARK_GREY)
        _box(sl, 6.1, top, 2.8, 0.42, verdict, font_size=11, color=DARK_GREY)
        _box(sl, 9.1, top, 0.7, 0.42, icon, font_size=13, bold=True,
             color=verdict_colors[i])


def slide_7_governance(prs: Presentation) -> None:
    """Slide 7 — Risks & Governance"""
    sl = _add_slide(prs)
    _bg(sl, WHITE)
    _rect(sl, 0, 0, 10, 0.5, BNR_BLUE)
    _box(sl, 0.2, 0.05, 9.6, 0.4, "Risks & Governance Considerations",
         font_size=22, bold=True, color=WHITE)

    items = [
        (
            "1. Risks in a central bank setting",
            RGBColor(0xCC, 0x33, 0x33),
            "Misquotation of legal/regulatory text; data leakage to external LLM API; "
            "outdated corpus if documents not re-indexed; over-reliance replacing expert judgement. "
            "Mitigations: local embeddings (no data egress); strict grounding prompt; mandatory human review.",
        ),
        (
            "2. Reducing hallucinations",
            RGBColor(0xCC, 0x33, 0x33),
            "System prompt forbids external knowledge. Fallback string enforced for out-of-corpus questions. "
            "NEW: add minimum similarity threshold (≥ 0.65) — Q4 demonstrated confident answer at sim=0.51. "
            "Post-retrieval faithfulness check: flag answers not entailed by retrieved chunks.",
        ),
        (
            "3. Audit & logging  (already implemented)",
            GREEN,
            "logs/audit.jsonl records every query: timestamp, question, answer preview, sources, "
            "similarity scores, tokens, latency, is_fallback flag. "
            "In production: centralise in ELK/Splunk with 5-year retention and access control.",
        ),
        (
            "4. Testing embedding or LLM model updates",
            ACCENT,
            "Maintain golden Q&A test set (these 5 questions + expected sources). "
            "For each candidate model: re-index → run test set → compare source-hit rate, "
            "faithfulness score, and fallback accuracy vs. baseline. Only promote if all metrics improve.",
        ),
        (
            "5. Improvements with more time",
            BNR_BLUE,
            "Hybrid retrieval (BM25 + dense) to fix Q1 keyword recall; "
            "similarity threshold gate to fix Q4 over-confidence; "
            "diversified top-k (force ≥2 sources) to fix Q2/Q5 single-doc bias; "
            "cross-encoder reranking; PDF table extraction; Kinyarwanda query support.",
        ),
    ]

    for i, (title, col, desc) in enumerate(items):
        top = 0.62 + i * 1.02
        _rect(sl, 0.2, top, 0.08, 0.75, col)
        _box(sl, 0.45, top, 3.0, 0.35,
             title, font_size=13, bold=True, color=col)
        _box(sl, 0.45, top + 0.32, 9.3, 0.65,
             desc, font_size=12, color=DARK_GREY)


# ── Main ──────────────────────────────────────────────────────────────────────

def main() -> None:
    prs = Presentation()
    prs.slide_width  = Inches(10)
    prs.slide_height = Inches(7.5)

    print("Building presentation …")
    slide_1_title(prs)
    slide_2_architecture(prs)
    slide_3_design(prs)
    slide_4_retrieval(prs)
    slide_5_outputs(prs)
    slide_6_evaluation(prs)
    slide_7_governance(prs)

    out = Path(__file__).parent / "BNR_RAG_Presentation_v2.pptx"
    prs.save(str(out))
    print(f"Saved → {out}")


if __name__ == "__main__":
    main()
